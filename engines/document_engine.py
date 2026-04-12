import os
import uuid
import io
import copy
from typing import List, Dict, Tuple, Optional, Union, Any
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptxRGBColor
import openpyxl
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.utils import get_column_letter
from openpyxl.chart import BarChart, LineChart, PieChart, Reference
import csv
import json
from datetime import datetime

class DocumentEngine:
    """
    Advanced document processing engine for:
    - Microsoft Word (DOCX)
    - Microsoft Excel (XLSX)
    - Microsoft PowerPoint (PPTX)
    
    Features:
    - Creation, Reading, Writing
    - Formatting and Styling
    - Tables and Charts
    - Images and Media
    - Templates
    - Conversion
    - Data Extraction
    - Merging and Splitting
    """
    
    def __init__(self):
        self.temp_dir = "temp_docs"
        os.makedirs(self.temp_dir, exist_ok=True)
    
    # =================== WORD DOCUMENT OPERATIONS ===================
    
    def create_word_document(self, output_path: str, title: str = "", 
                            content: List[Dict] = None) -> str:
        """
        Create Word document with structured content
        
        Args:
            output_path: Output file path
            title: Document title
            content: List of content blocks, each with 'type' and data
                    Types: 'heading', 'paragraph', 'table', 'image', 'list'
        
        Example:
            content = [
                {'type': 'heading', 'text': 'Chapter 1', 'level': 1},
                {'type': 'paragraph', 'text': 'This is a paragraph.'},
                {'type': 'table', 'data': [[...], [...]]},
                {'type': 'image', 'path': 'image.jpg', 'width': 4}
            ]
        """
        doc = Document()
        
        # Add title
        if title:
            doc.add_heading(title, 0)
        
        # Add content blocks
        if content:
            for block in content:
                block_type = block.get('type', 'paragraph')
                
                if block_type == 'heading':
                    level = block.get('level', 1)
                    text = block.get('text', '')
                    doc.add_heading(text, level)
                
                elif block_type == 'paragraph':
                    text = block.get('text', '')
                    style = block.get('style', None)
                    p = doc.add_paragraph(text, style=style)
                    
                    # Apply formatting
                    if block.get('bold'):
                        p.runs[0].bold = True
                    if block.get('italic'):
                        p.runs[0].italic = True
                    if block.get('font_size'):
                        p.runs[0].font.size = Pt(block['font_size'])
                    if block.get('color'):
                        r, g, b = block['color']
                        p.runs[0].font.color.rgb = RGBColor(r, g, b)
                    if block.get('alignment'):
                        align_map = {
                            'left': WD_ALIGN_PARAGRAPH.LEFT,
                            'center': WD_ALIGN_PARAGRAPH.CENTER,
                            'right': WD_ALIGN_PARAGRAPH.RIGHT,
                            'justify': WD_ALIGN_PARAGRAPH.JUSTIFY
                        }
                        p.alignment = align_map.get(block['alignment'])
                
                elif block_type == 'table':
                    data = block.get('data', [])
                    if data:
                        table = doc.add_table(rows=len(data), cols=len(data[0]))
                        table.style = block.get('style', 'Table Grid')
                        
                        for i, row in enumerate(data):
                            for j, cell_value in enumerate(row):
                                table.rows[i].cells[j].text = str(cell_value)
                
                elif block_type == 'image':
                    image_path = block.get('path')
                    width = block.get('width', 4)  # inches
                    if image_path and os.path.exists(image_path):
                        doc.add_picture(image_path, width=Inches(width))
                
                elif block_type == 'list':
                    items = block.get('items', [])
                    for item in items:
                        doc.add_paragraph(item, style='List Bullet')
                
                elif block_type == 'numbered_list':
                    items = block.get('items', [])
                    for item in items:
                        doc.add_paragraph(item, style='List Number')
        
        doc.save(output_path)
        return output_path
    
    def read_word_document(self, docx_path: str) -> Dict:
        """
        Read Word document and extract all content
        
        Returns:
            Dict with 'paragraphs', 'tables', 'images', 'metadata'
        """
        doc = Document(docx_path)
        
        result = {
            'paragraphs': [],
            'tables': [],
            'images': [],
            'metadata': {}
        }
        
        # Extract paragraphs
        for para in doc.paragraphs:
            result['paragraphs'].append({
                'text': para.text,
                'style': para.style.name if para.style else None
            })
        
        # Extract tables
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            result['tables'].append(table_data)
        
        # Extract images
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                result['images'].append({
                    'type': rel.target_part.content_type,
                    'size': len(rel.target_part.blob)
                })
        
        # Extract metadata
        core_props = doc.core_properties
        result['metadata'] = {
            'title': core_props.title or '',
            'author': core_props.author or '',
            'subject': core_props.subject or '',
            'created': str(core_props.created) if core_props.created else '',
            'modified': str(core_props.modified) if core_props.modified else ''
        }
        
        return result
    
    def merge_word_documents(self, docx_paths: List[str], output_path: str,
                            add_page_breaks: bool = True) -> str:
        """
        Merge multiple Word documents
        
        Args:
            docx_paths: List of document paths
            output_path: Output file path
            add_page_breaks: Add page break between documents
        """
        merged_doc = Document()
        
        for i, docx_path in enumerate(docx_paths):
            doc = Document(docx_path)
            
            # Add page break between documents
            if i > 0 and add_page_breaks:
                merged_doc.add_page_break()
            
            # Copy all elements
            for element in doc.element.body:
                merged_doc.element.body.append(element)
        
        merged_doc.save(output_path)
        return output_path
    
    def add_table_to_word(self, docx_path: str, output_path: str,
                         data: List[List[Any]], position: int = -1,
                         style: str = 'Table Grid',
                         header_row: bool = True) -> str:
        """
        Add table to existing Word document
        
        Args:
            docx_path: Input document path
            output_path: Output document path
            data: 2D list of table data
            position: Paragraph index to insert after (-1 for end)
            style: Table style name
            header_row: First row is header
        """
        doc = Document(docx_path)
        
        # Create table
        table = doc.add_table(rows=len(data), cols=len(data[0]))
        table.style = style
        
        # Populate table
        for i, row in enumerate(data):
            for j, cell_value in enumerate(row):
                cell = table.rows[i].cells[j]
                cell.text = str(cell_value)
                
                # Format header row
                if i == 0 and header_row:
                    cell.paragraphs[0].runs[0].bold = True
        
        doc.save(output_path)
        return output_path
    
    def replace_text_in_word(self, docx_path: str, output_path: str,
                            replacements: Dict[str, str]) -> str:
        """
        Find and replace text in Word document
        
        Args:
            docx_path: Input document path
            output_path: Output document path
            replacements: Dict mapping old text to new text
        """
        doc = Document(docx_path)
        
        for paragraph in doc.paragraphs:
            for old_text, new_text in replacements.items():
                if old_text in paragraph.text:
                    # Replace in each run to preserve formatting
                    for run in paragraph.runs:
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)
        
        # Also replace in tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for old_text, new_text in replacements.items():
                        if old_text in cell.text:
                            cell.text = cell.text.replace(old_text, new_text)
        
        doc.save(output_path)
        return output_path
    
    def set_word_metadata(self, docx_path: str, output_path: str,
                         metadata: Dict) -> str:
        """
        Set Word document metadata
        
        Args:
            metadata: Dict with title, author, subject, keywords, etc.
        """
        doc = Document(docx_path)
        
        core_props = doc.core_properties
        
        if 'title' in metadata:
            core_props.title = metadata['title']
        if 'author' in metadata:
            core_props.author = metadata['author']
        if 'subject' in metadata:
            core_props.subject = metadata['subject']
        if 'keywords' in metadata:
            core_props.keywords = metadata['keywords']
        
        doc.save(output_path)
        return output_path
    
    # =================== EXCEL OPERATIONS ===================
    
    def create_excel_workbook(self, output_path: str, 
                             sheets: Dict[str, List[List[Any]]] = None,
                             formatting: Dict = None) -> str:
        """
        Create Excel workbook with multiple sheets
        
        Args:
            output_path: Output file path
            sheets: Dict mapping sheet names to 2D data arrays
            formatting: Optional formatting settings
        
        Example:
            sheets = {
                'Sales': [['Product', 'Q1', 'Q2'], ['A', 100, 150], ...],
                'Expenses': [[...], [...]]
            }
        """
        wb = openpyxl.Workbook()
        
        # Remove default sheet
        if 'Sheet' in wb.sheetnames:
            wb.remove(wb['Sheet'])
        
        if sheets:
            for sheet_name, data in sheets.items():
                ws = wb.create_sheet(sheet_name)
                
                # Populate data
                for row_idx, row in enumerate(data, 1):
                    for col_idx, value in enumerate(row, 1):
                        ws.cell(row=row_idx, column=col_idx, value=value)
                
                # Apply formatting if provided
                if formatting and sheet_name in formatting:
                    self._apply_excel_formatting(ws, formatting[sheet_name])
        else:
            # Create empty sheet
            wb.create_sheet('Sheet1')
        
        wb.save(output_path)
        return output_path
    
    def read_excel_workbook(self, xlsx_path: str, 
                           sheet_names: Optional[List[str]] = None) -> Dict:
        """
        Read Excel workbook
        
        Args:
            xlsx_path: Excel file path
            sheet_names: Specific sheets to read (None for all)
        
        Returns:
            Dict mapping sheet names to data arrays
        """
        wb = openpyxl.load_workbook(xlsx_path)
        result = {}
        
        sheets_to_read = sheet_names if sheet_names else wb.sheetnames
        
        for sheet_name in sheets_to_read:
            if sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                data = []
                
                for row in ws.iter_rows(values_only=True):
                    data.append(list(row))
                
                result[sheet_name] = data
        
        return result
    
    def merge_excel_workbooks(self, xlsx_paths: List[str], output_path: str,
                             mode: str = 'sheets') -> str:
        """
        Merge multiple Excel workbooks
        
        Args:
            xlsx_paths: List of Excel file paths
            output_path: Output file path
            mode: 'sheets' (each file as separate sheets) or 
                  'rows' (combine all into one sheet)
        """
        wb_merged = openpyxl.Workbook()
        wb_merged.remove(wb_merged.active)  # Remove default sheet
        
        if mode == 'sheets':
            # Each file becomes separate sheets
            for i, xlsx_path in enumerate(xlsx_paths):
                wb = openpyxl.load_workbook(xlsx_path)
                
                for sheet_name in wb.sheetnames:
                    ws_source = wb[sheet_name]
                    
                    # Create new sheet with unique name
                    new_name = f"File{i+1}_{sheet_name}"
                    ws_dest = wb_merged.create_sheet(new_name)
                    
                    # Copy data
                    for row in ws_source.iter_rows():
                        for cell in row:
                            ws_dest[cell.coordinate].value = cell.value
        
        else:  # mode == 'rows'
            # Combine all data into one sheet
            ws_merged = wb_merged.create_sheet('Merged Data')
            current_row = 1
            
            for i, xlsx_path in enumerate(xlsx_paths):
                wb = openpyxl.load_workbook(xlsx_path)
                ws = wb.active
                
                for row in ws.iter_rows(values_only=True):
                    # Skip header rows except for first file
                    if i > 0 and current_row == 1:
                        continue
                    
                    ws_merged.append(row)
                    current_row += 1
        
        wb_merged.save(output_path)
        return output_path
    
    def add_excel_chart(self, xlsx_path: str, output_path: str,
                       sheet_name: str, chart_type: str,
                       data_range: str, chart_position: str = 'E2') -> str:
        """
        Add chart to Excel worksheet
        
        Args:
            xlsx_path: Input file path
            output_path: Output file path
            sheet_name: Sheet to add chart to
            chart_type: 'bar', 'line', 'pie'
            data_range: Data range (e.g., 'A1:B10')
            chart_position: Cell position for chart
        """
        wb = openpyxl.load_workbook(xlsx_path)
        ws = wb[sheet_name]
        
        # Create chart based on type
        if chart_type == 'bar':
            chart = BarChart()
        elif chart_type == 'line':
            chart = LineChart()
        elif chart_type == 'pie':
            chart = PieChart()
        else:
            raise ValueError(f"Unknown chart type: {chart_type}")
        
        # Parse data range
        data = Reference(ws, range_string=data_range)
        chart.add_data(data, titles_from_data=True)
        
        # Add chart to worksheet
        ws.add_chart(chart, chart_position)
        
        wb.save(output_path)
        return output_path
    
    def format_excel_cells(self, xlsx_path: str, output_path: str,
                          sheet_name: str, cell_range: str,
                          formatting: Dict) -> str:
        """
        Format Excel cells
        
        Args:
            xlsx_path: Input file path
            output_path: Output file path
            sheet_name: Sheet name
            cell_range: Cell range (e.g., 'A1:D10')
            formatting: Dict with formatting options
        
        Formatting options:
            - font_name, font_size, font_color, bold, italic
            - bg_color
            - border
            - alignment (left, center, right)
            - number_format
        """
        wb = openpyxl.load_workbook(xlsx_path)
        ws = wb[sheet_name]
        
        # Parse range
        cells = ws[cell_range]
        
        for row in cells:
            for cell in (row if isinstance(row, tuple) else [row]):
                # Font
                font_kwargs = {}
                if 'font_name' in formatting:
                    font_kwargs['name'] = formatting['font_name']
                if 'font_size' in formatting:
                    font_kwargs['size'] = formatting['font_size']
                if 'bold' in formatting:
                    font_kwargs['bold'] = formatting['bold']
                if 'italic' in formatting:
                    font_kwargs['italic'] = formatting['italic']
                if 'font_color' in formatting:
                    font_kwargs['color'] = formatting['font_color']
                
                if font_kwargs:
                    cell.font = Font(**font_kwargs)
                
                # Background color
                if 'bg_color' in formatting:
                    cell.fill = PatternFill(start_color=formatting['bg_color'],
                                          end_color=formatting['bg_color'],
                                          fill_type='solid')
                
                # Alignment
                if 'alignment' in formatting:
                    align_map = {
                        'left': 'left',
                        'center': 'center',
                        'right': 'right'
                    }
                    cell.alignment = Alignment(horizontal=align_map.get(formatting['alignment']))
                
                # Number format
                if 'number_format' in formatting:
                    cell.number_format = formatting['number_format']
                
                # Border
                if 'border' in formatting and formatting['border']:
                    side = Side(style='thin', color='000000')
                    cell.border = Border(left=side, right=side, top=side, bottom=side)
        
        wb.save(output_path)
        return output_path
    
    def excel_to_csv(self, xlsx_path: str, output_dir: str,
                    sheet_names: Optional[List[str]] = None) -> List[str]:
        """
        Convert Excel sheets to CSV files
        
        Args:
            xlsx_path: Excel file path
            output_dir: Output directory
            sheet_names: Specific sheets to convert (None for all)
        
        Returns:
            List of created CSV file paths
        """
        os.makedirs(output_dir, exist_ok=True)
        wb = openpyxl.load_workbook(xlsx_path)
        output_files = []
        
        sheets_to_convert = sheet_names if sheet_names else wb.sheetnames
        
        for sheet_name in sheets_to_convert:
            if sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                
                csv_filename = f"{sheet_name}.csv"
                csv_path = os.path.join(output_dir, csv_filename)
                
                with open(csv_path, 'w', newline='', encoding='utf-8') as f:
                    writer = csv.writer(f)
                    
                    for row in ws.iter_rows(values_only=True):
                        writer.writerow(row)
                
                output_files.append(csv_path)
        
        return output_files
    
    def csv_to_excel(self, csv_paths: List[str], output_path: str) -> str:
        """
        Convert CSV files to Excel workbook
        
        Args:
            csv_paths: List of CSV file paths
            output_path: Output Excel file path
        """
        wb = openpyxl.Workbook()
        wb.remove(wb.active)  # Remove default sheet
        
        for csv_path in csv_paths:
            sheet_name = os.path.splitext(os.path.basename(csv_path))[0]
            ws = wb.create_sheet(sheet_name)
            
            with open(csv_path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                
                for row in reader:
                    ws.append(row)
        
        wb.save(output_path)
        return output_path
    
    def _apply_excel_formatting(self, ws, formatting: Dict):
        """Apply formatting to worksheet"""
        # Header formatting
        if 'header_row' in formatting:
            header_row = formatting['header_row']
            for cell in ws[header_row]:
                cell.font = Font(bold=True)
                cell.fill = PatternFill(start_color='CCCCCC', 
                                      end_color='CCCCCC',
                                      fill_type='solid')
        
        # Auto-size columns
        if formatting.get('auto_width', False):
            for column in ws.columns:
                max_length = 0
                column_letter = get_column_letter(column[0].column)
                
                for cell in column:
                    if cell.value:
                        max_length = max(max_length, len(str(cell.value)))
                
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column_letter].width = adjusted_width
    
    # =================== POWERPOINT OPERATIONS ===================
    
    def create_presentation(self, output_path: str, title: str = "",
                          slides: List[Dict] = None) -> str:
        """
        Create PowerPoint presentation
        
        Args:
            output_path: Output file path
            title: Presentation title
            slides: List of slide definitions
        
        Slide types:
            - title_slide: {'type': 'title_slide', 'title': '...', 'subtitle': '...'}
            - title_content: {'type': 'title_content', 'title': '...', 'content': [...]}
            - blank: {'type': 'blank'}
        """
        prs = Presentation()
        
        if slides:
            for slide_def in slides:
                slide_type = slide_def.get('type', 'title_content')
                
                if slide_type == 'title_slide':
                    self._add_title_slide(prs, slide_def)
                elif slide_type == 'title_content':
                    self._add_title_content_slide(prs, slide_def)
                elif slide_type == 'blank':
                    prs.slides.add_slide(prs.slide_layouts[6])
                elif slide_type == 'image':
                    self._add_image_slide(prs, slide_def)
        else:
            # Add title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            
            title_shape.text = title if title else "Presentation Title"
            subtitle_shape.text = "Created with DocumentEngine"
        
        prs.save(output_path)
        return output_path
    
    def _add_title_slide(self, prs, slide_def):
        """Add title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = slide_def.get('title', '')
        subtitle.text = slide_def.get('subtitle', '')
    
    def _add_title_content_slide(self, prs, slide_def):
        """Add title and content slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = slide_def.get('title', '')
        
        # Add content
        content = slide_def.get('content', [])
        if content:
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            
            for i, item in enumerate(content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = item
                p.level = 0
    
    def _add_image_slide(self, prs, slide_def):
        """Add image slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        image_path = slide_def.get('image_path')
        if image_path and os.path.exists(image_path):
            left = PptxInches(1)
            top = PptxInches(1)
            width = PptxInches(8)
            
            slide.shapes.add_picture(image_path, left, top, width=width)
        
        # Add title if provided
        if slide_def.get('title'):
            left = PptxInches(0.5)
            top = PptxInches(0.5)
            width = PptxInches(9)
            height = PptxInches(1)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = slide_def['title']
            
            # Format title
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = PptxPt(32)
                paragraph.font.bold = True
    
    def read_presentation(self, pptx_path: str) -> Dict:
        """
        Read PowerPoint presentation
        
        Returns:
            Dict with slide content and metadata
        """
        prs = Presentation(pptx_path)
        
        result = {
            'slides': [],
            'metadata': {
                'slide_count': len(prs.slides),
                'slide_width': prs.slide_width,
                'slide_height': prs.slide_height
            }
        }
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_data = {
                'slide_number': slide_num,
                'shapes': []
            }
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_data['shapes'].append({
                        'type': 'text',
                        'text': shape.text
                    })
                elif hasattr(shape, "image"):
                    slide_data['shapes'].append({
                        'type': 'image',
                        'size': len(shape.image.blob)
                    })
            
            result['slides'].append(slide_data)
        
        return result
    
    def merge_presentations(self, pptx_paths: List[str], output_path: str) -> str:
        """
        Merge multiple PowerPoint presentations
        
        Args:
            pptx_paths: List of presentation paths
            output_path: Output file path
        """
        # Start with the first presentation
        merged_prs = Presentation(pptx_paths[0])
        
        # Add slides from remaining presentations
        for pptx_path in pptx_paths[1:]:
            prs = Presentation(pptx_path)
            
            for slide in prs.slides:
                # Get the slide layout
                slide_layout = slide.slide_layout
                
                # Try to find matching layout in merged presentation
                try:
                    # Use blank layout as fallback
                    new_slide = merged_prs.slides.add_slide(merged_prs.slide_layouts[6])
                    
                    # Copy all shapes from original slide
                    for shape in slide.shapes:
                        # Create a deep copy of the shape element
                        el = shape.element
                        newel = copy.deepcopy(el)
                        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
                
                except Exception as e:
                    # If copying fails, try simpler approach
                    # Just copy text content
                    new_slide = merged_prs.slides.add_slide(merged_prs.slide_layouts[6])
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            # Add text box with the content
                            left = shape.left
                            top = shape.top
                            width = shape.width
                            height = shape.height
                            
                            txBox = new_slide.shapes.add_textbox(left, top, width, height)
                            tf = txBox.text_frame
                            tf.text = shape.text
        
        merged_prs.save(output_path)
        return output_path
    
    # =================== CONVERSION OPERATIONS ===================
    
    def word_to_text(self, docx_path: str, output_path: str) -> str:
        """Convert Word document to plain text"""
        doc = Document(docx_path)
        
        text = "\n\n".join([para.text for para in doc.paragraphs])
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        
        return output_path
    
    def excel_to_json(self, xlsx_path: str, output_path: str,
                     sheet_name: Optional[str] = None) -> str:
        """
        Convert Excel to JSON
        
        Args:
            xlsx_path: Excel file path
            output_path: JSON output path
            sheet_name: Specific sheet (None for active sheet)
        """
        wb = openpyxl.load_workbook(xlsx_path)
        
        if sheet_name:
            ws = wb[sheet_name]
        else:
            ws = wb.active
        
        # Get headers from first row
        headers = [cell.value for cell in next(ws.iter_rows())]
        
        # Convert to list of dicts
        data = []
        for row in ws.iter_rows(min_row=2, values_only=True):
            row_dict = {headers[i]: value for i, value in enumerate(row)}
            data.append(row_dict)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=2, default=str)
        
        return output_path
    
    # =================== UTILITY METHODS ===================
    
    def cleanup_temp_files(self):
        """Clean up temporary files"""
        import shutil
        if os.path.exists(self.temp_dir):
            shutil.rmtree(self.temp_dir)
            os.makedirs(self.temp_dir, exist_ok=True)