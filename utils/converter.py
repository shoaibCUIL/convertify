import os
import uuid
from PIL import Image
import PyPDF2
from pdf2image import convert_from_path
from docx import Document
from pptx import Presentation
import openpyxl
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
import img2pdf

class UniversalConverter:
    """Convert between various file formats"""
    
    def __init__(self):
        self.supported_conversions = {
            'pdf': ['docx', 'jpg', 'png', 'txt', 'html'],
            'docx': ['pdf', 'txt', 'html'],
            'xlsx': ['pdf', 'csv', 'html'],
            'pptx': ['pdf', 'jpg', 'png'],
            'jpg': ['pdf', 'png', 'bmp', 'gif', 'webp'],
            'jpeg': ['pdf', 'png', 'bmp', 'gif', 'webp'],
            'png': ['pdf', 'jpg', 'bmp', 'gif', 'webp'],
            'gif': ['pdf', 'jpg', 'png', 'bmp'],
            'bmp': ['pdf', 'jpg', 'png', 'gif'],
            'txt': ['pdf', 'docx', 'html'],
            'html': ['pdf', 'docx', 'txt'],
            'csv': ['xlsx', 'pdf']
        }
    
    def convert(self, input_file, target_format, output_folder):
        """Main conversion method"""
        source_format = os.path.splitext(input_file)[1][1:].lower()
        target_format = target_format.lower()
        
        if source_format not in self.supported_conversions:
            raise ValueError(f"Unsupported source format: {source_format}")
        
        if target_format not in self.supported_conversions.get(source_format, []):
            raise ValueError(f"Cannot convert {source_format} to {target_format}")
        
        # Generate output filename
        output_filename = f"{uuid.uuid4()}.{target_format}"
        output_path = os.path.join(output_folder, output_filename)
        
        # Route to appropriate converter
        converter_method = f"_{source_format}_to_{target_format}"
        if hasattr(self, converter_method):
            getattr(self, converter_method)(input_file, output_path)
        else:
            raise NotImplementedError(f"Conversion {source_format} -> {target_format} not implemented")
        
        return output_path
    
    # =================== PDF CONVERSIONS ===================
    def _pdf_to_jpg(self, input_file, output_path):
        """Convert PDF to JPG (first page)"""
        images = convert_from_path(input_file, dpi=300)
        images[0].save(output_path, 'JPEG', quality=95)
    
    def _pdf_to_png(self, input_file, output_path):
        """Convert PDF to PNG (first page)"""
        images = convert_from_path(input_file, dpi=300)
        images[0].save(output_path, 'PNG')
    
    def _pdf_to_txt(self, input_file, output_path):
        """Extract text from PDF"""
        text = ""
        with open(input_file, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n\n"
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
    
    def _pdf_to_docx(self, input_file, output_path):
        """Convert PDF to DOCX"""
        # Extract text and create DOCX
        text = ""
        with open(input_file, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n\n"
        
        doc = Document()
        for paragraph in text.split('\n\n'):
            if paragraph.strip():
                doc.add_paragraph(paragraph.strip())
        doc.save(output_path)
    
    def _pdf_to_html(self, input_file, output_path):
        """Convert PDF to HTML"""
        text = ""
        with open(input_file, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n\n"
        
        html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>Converted PDF</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }}
    </style>
</head>
<body>
    <div>{text.replace(chr(10), '<br>')}</div>
</body>
</html>"""
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html)
    
    # =================== IMAGE CONVERSIONS ===================
    def _jpg_to_pdf(self, input_file, output_path):
        """Convert JPG to PDF"""
        with open(output_path, 'wb') as f:
            f.write(img2pdf.convert(input_file))
    
    def _png_to_pdf(self, input_file, output_path):
        """Convert PNG to PDF"""
        with open(output_path, 'wb') as f:
            f.write(img2pdf.convert(input_file))
    
    def _jpg_to_png(self, input_file, output_path):
        """Convert JPG to PNG"""
        img = Image.open(input_file)
        img.save(output_path, 'PNG')
    
    def _png_to_jpg(self, input_file, output_path):
        """Convert PNG to JPG"""
        img = Image.open(input_file)
        if img.mode in ('RGBA', 'LA', 'P'):
            # Create white background
            background = Image.new('RGB', img.size, (255, 255, 255))
            if img.mode == 'P':
                img = img.convert('RGBA')
            background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
            img = background
        img.save(output_path, 'JPEG', quality=95)
    
    def _jpg_to_bmp(self, input_file, output_path):
        """Convert JPG to BMP"""
        img = Image.open(input_file)
        img.save(output_path, 'BMP')
    
    def _png_to_bmp(self, input_file, output_path):
        """Convert PNG to BMP"""
        img = Image.open(input_file)
        img.save(output_path, 'BMP')
    
    def _jpg_to_gif(self, input_file, output_path):
        """Convert JPG to GIF"""
        img = Image.open(input_file)
        img.save(output_path, 'GIF')
    
    def _png_to_gif(self, input_file, output_path):
        """Convert PNG to GIF"""
        img = Image.open(input_file)
        img.save(output_path, 'GIF')
    
    def _jpg_to_webp(self, input_file, output_path):
        """Convert JPG to WebP"""
        img = Image.open(input_file)
        img.save(output_path, 'WEBP', quality=90)
    
    def _png_to_webp(self, input_file, output_path):
        """Convert PNG to WebP"""
        img = Image.open(input_file)
        img.save(output_path, 'WEBP', quality=90)
    
    def _gif_to_jpg(self, input_file, output_path):
        """Convert GIF to JPG"""
        img = Image.open(input_file)
        img = img.convert('RGB')
        img.save(output_path, 'JPEG', quality=95)
    
    def _gif_to_png(self, input_file, output_path):
        """Convert GIF to PNG"""
        img = Image.open(input_file)
        img.save(output_path, 'PNG')
    
    def _bmp_to_jpg(self, input_file, output_path):
        """Convert BMP to JPG"""
        img = Image.open(input_file)
        img.save(output_path, 'JPEG', quality=95)
    
    def _bmp_to_png(self, input_file, output_path):
        """Convert BMP to PNG"""
        img = Image.open(input_file)
        img.save(output_path, 'PNG')
    
    # =================== DOCUMENT CONVERSIONS ===================
    def _docx_to_pdf(self, input_file, output_path):
        """Convert DOCX to PDF"""
        from docx2pdf import convert
        convert(input_file, output_path)
    
    def _docx_to_txt(self, input_file, output_path):
        """Convert DOCX to TXT"""
        doc = Document(input_file)
        text = "\n\n".join([para.text for para in doc.paragraphs])
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
    
    def _docx_to_html(self, input_file, output_path):
        """Convert DOCX to HTML"""
        doc = Document(input_file)
        html = """<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>Converted Document</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }
    </style>
</head>
<body>
"""
        for para in doc.paragraphs:
            html += f"    <p>{para.text}</p>\n"
        html += "</body>\n</html>"
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html)
    
    def _txt_to_pdf(self, input_file, output_path):
        """Convert TXT to PDF"""
        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter
        
        with open(input_file, 'r', encoding='utf-8') as f:
            text = f.read()
        
        y = height - 50
        for line in text.split('\n'):
            if y < 50:
                c.showPage()
                y = height - 50
            c.drawString(50, y, line[:100])  # Limit line length
            y -= 15
        
        c.save()
    
    def _txt_to_docx(self, input_file, output_path):
        """Convert TXT to DOCX"""
        with open(input_file, 'r', encoding='utf-8') as f:
            text = f.read()
        
        doc = Document()
        for paragraph in text.split('\n\n'):
            if paragraph.strip():
                doc.add_paragraph(paragraph.strip())
        doc.save(output_path)
    
    # =================== SPREADSHEET CONVERSIONS ===================
    def _xlsx_to_pdf(self, input_file, output_path):
        """Convert XLSX to PDF"""
        wb = openpyxl.load_workbook(input_file)
        ws = wb.active
        
        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter
        
        y = height - 50
        for row in ws.iter_rows(values_only=True):
            if y < 50:
                c.showPage()
                y = height - 50
            text = " | ".join([str(cell) if cell is not None else "" for cell in row])
            c.drawString(50, y, text[:100])
            y -= 15
        
        c.save()
    
    def _xlsx_to_csv(self, input_file, output_path):
        """Convert XLSX to CSV"""
        wb = openpyxl.load_workbook(input_file)
        ws = wb.active
        
        import csv
        with open(output_path, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            for row in ws.iter_rows(values_only=True):
                writer.writerow(row)
    
    def _csv_to_xlsx(self, input_file, output_path):
        """Convert CSV to XLSX"""
        import csv
        wb = openpyxl.Workbook()
        ws = wb.active
        
        with open(input_file, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            for row in reader:
                ws.append(row)
        
        wb.save(output_path)
    
    # =================== PRESENTATION CONVERSIONS ===================
    def _pptx_to_pdf(self, input_file, output_path):
        """Convert PPTX to PDF (simplified)"""
        prs = Presentation(input_file)
        
        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter
        
        for slide_num, slide in enumerate(prs.slides, 1):
            c.drawString(50, height - 50, f"Slide {slide_num}")
            y = height - 80
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    if y < 50:
                        break
                    c.drawString(50, y, shape.text[:100])
                    y -= 15
            
            c.showPage()
        
        c.save()
    
    def _pptx_to_jpg(self, input_file, output_path):
        """Convert PPTX to JPG (first slide preview)"""
        # This is a simplified version - full implementation would need more libs
        prs = Presentation(input_file)
        
        # Create a simple image with text
        img = Image.new('RGB', (800, 600), color='white')
        # In production, you'd use a proper rendering library
        img.save(output_path, 'JPEG', quality=95)